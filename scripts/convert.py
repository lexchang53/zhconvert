#!/usr/bin/env python
# -*- coding: utf-8 -*-

import os
import sys
import argparse
import subprocess

def install_package(package_name, install_name=None):
    """如果缺少套件，自動在背景以 pip 安裝"""
    install_name = install_name or package_name
    print(f"[*] 偵測到缺少必要套件 {package_name}，正在嘗試自動安裝 {install_name}...")
    try:
        # 使用目前的 Python 直譯器執行 pip 安裝，並將輸出導向 null 以保持畫面乾淨
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", install_name],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL
        )
        print(f"[+] 成功安裝 {package_name}！")
    except Exception as e:
        print(f"[-] 自動安裝 {package_name} 失敗: {e}。請手動執行 'pip install {install_name}'")
        sys.exit(1)

# 動態檢測與導入必要套件
try:
    import requests
except ImportError:
    install_package("requests")
    import requests


API_URL = "https://api.zhconvert.org/convert"
MAX_CHUNK_SIZE = 5000  # 繁化姬單次 API 建議最大發送字數

def _call_api(text, converter):
    """實際呼叫繁化姬 API 的底層函數"""
    payload = {
        "text": text,
        "converter": converter,  # 'Taiwan' 或 'Traditional'
        "outputFormat": "json"
    }
    try:
        response = requests.post(API_URL, data=payload, timeout=30)
        if response.status_code == 200:
            res_json = response.json()
            if "data" in res_json and "text" in res_json["data"]:
                return res_json["data"]["text"]
            else:
                print(f"[!] API 回傳格式異常: {res_json}")
        else:
            print(f"[!] API 請求失敗，狀態碼: {response.status_code}")
    except Exception as e:
        print(f"[!] 呼叫 API 時發生連線錯誤: {e}")
    return text

def convert_text(text, converter="Taiwan"):
    """
    呼叫繁化姬 API 進行繁體化與台灣化。
    支援大文本分段處理，避免超出 API 負載限制。
    """
    if not text or not text.strip():
        return text

    if len(text) <= MAX_CHUNK_SIZE:
        return _call_api(text, converter)

    # 超過限制則按行分段，避免切斷單個詞彙
    lines = text.split("\n")
    converted_lines = []
    current_chunk = []
    current_len = 0

    for line in lines:
        line_len = len(line) + 1  # 算入換行符
        if current_len + line_len > MAX_CHUNK_SIZE:
            if current_chunk:
                chunk_text = "\n".join(current_chunk)
                converted_lines.append(_call_api(chunk_text, converter))
                current_chunk = []
                current_len = 0
            # 如果單行就超過最大限制，強制依字數切分
            if line_len > MAX_CHUNK_SIZE:
                for i in range(0, len(line), MAX_CHUNK_SIZE):
                    sub_chunk = line[i:i+MAX_CHUNK_SIZE]
                    converted_lines.append(_call_api(sub_chunk, converter))
                continue
        current_chunk.append(line)
        current_len += line_len

    if current_chunk:
        chunk_text = "\n".join(current_chunk)
        converted_lines.append(_call_api(chunk_text, converter))

    return "\n".join(converted_lines)


# ==================== 檔案格式處理邏輯 ====================

def convert_text_file(file_path, output_path, converter):
    """處理純文字或 Markdown 檔案"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    
    converted = convert_text(content, converter)
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(converted)
    print(f"[OK] 已成功轉換純文字/Markdown: {os.path.basename(file_path)}")

def convert_docx_file(file_path, output_path, converter):
    """處理 Word 檔案（逐個 Run 替換以保留格式）"""
    try:
        import docx
    except ImportError:
        install_package("docx", "python-docx")
        import docx
        
    doc = docx.Document(file_path)
    
    # 1. 處理文件主體段落
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                run.text = convert_text(run.text, converter)
                
    # 2. 處理文件內的所有表格
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            run.text = convert_text(run.text, converter)
                            
    # 3. 處理頁首與頁尾
    for section in doc.sections:
        # 頁首
        for header in [section.header, section.first_page_header, section.even_page_header]:
            if header:
                for paragraph in header.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            run.text = convert_text(run.text, converter)
        # 頁尾
        for footer in [section.footer, section.first_page_footer, section.even_page_footer]:
            if footer:
                for paragraph in footer.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            run.text = convert_text(run.text, converter)
                            
    doc.save(output_path)
    print(f"[OK] 已成功轉換 Word 檔（格式已保留）: {os.path.basename(file_path)}")

def convert_xlsx_file(file_path, output_path, converter):
    """處理 Excel 檔案（保留儲存格樣式與公式）"""
    try:
        import openpyxl
    except ImportError:
        install_package("openpyxl")
        import openpyxl
        
    wb = openpyxl.load_workbook(file_path)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    # 跳過以 '=' 開頭的 Excel 公式，避免破壞函數結構
                    if cell.value.startswith("="):
                        continue
                    cell.value = convert_text(cell.value, converter)
    wb.save(output_path)
    print(f"[OK] 已成功轉換 Excel 檔（格式已保留）: {os.path.basename(file_path)}")

def convert_pptx_file(file_path, output_path, converter):
    """處理 PowerPoint 檔案（保留版面與投影片樣式）"""
    try:
        import pptx
    except ImportError:
        install_package("pptx", "python-pptx")
        import pptx
        
    prs = pptx.Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            # 處理文字框與自訂形狀中的文字
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            run.text = convert_text(run.text, converter)
            # 處理投影片中的表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text and run.text.strip():
                                    run.text = convert_text(run.text, converter)
    prs.save(output_path)
    print(f"[OK] 已成功轉換 PowerPoint 簡報（格式已保留）: {os.path.basename(file_path)}")


# ==================== 主程式排程與進入點 ====================

def process_file(file_path, converter, overwrite=False):
    """判斷副檔名並調用對應的處理函數，支援檔案名稱與內容的同步簡繁轉換"""
    if not os.path.exists(file_path):
        print(f"[-] 錯誤：找不到路徑 '{file_path}'")
        return False

    ext = os.path.splitext(file_path)[1].lower()
    
    dir_name, base_name = os.path.split(file_path)
    name, extension = os.path.splitext(base_name)
    
    # 將檔案名稱也進行簡繁與用語轉換
    try:
        new_name = convert_text(name, converter)
    except Exception:
        new_name = name
        
    backup_path = os.path.join(dir_name, f"{name}_old{extension}")
    output_path = os.path.join(dir_name, f"{new_name}{extension}")
    
    # 決定輸入與輸出路徑
    if overwrite:
        # 直接覆蓋模式
        input_path = file_path
    else:
        # 備份覆蓋模式：將原檔重新命名為 _old，轉換結果寫入新譯名檔
        # 安全防護：如果 _old 備份檔已存在，說明這可能是重複執行的轉檔。
        # 我們絕對不能覆蓋已有的 _old 簡體原檔，以保護您的原始數據！
        if os.path.exists(backup_path):
            print(f"[*] 備份檔 '{os.path.basename(backup_path)}' 已存在，略過備份以保護原始數據。")
            input_path = file_path
        else:
            try:
                # 第一次轉換，正常將簡體原檔備份為 _old
                os.replace(file_path, backup_path)
                input_path = backup_path
            except Exception as e:
                print(f"[-] 無法為檔案建立備份 {file_path}: {e}")
                return False

    try:
        if ext in [".md", ".txt"]:
            convert_text_file(input_path, output_path, converter)
        elif ext == ".docx":
            convert_docx_file(input_path, output_path, converter)
        elif ext == ".xlsx":
            convert_xlsx_file(input_path, output_path, converter)
        elif ext == ".pptx":
            convert_pptx_file(input_path, output_path, converter)
        else:
            print(f"[!] 略過不支援的檔案格式: {os.path.basename(file_path)}")
            # 若發生略過，且剛才有做重新命名，則將其復原為原名稱
            if not overwrite:
                os.replace(input_path, file_path)
            return False
            
        # 若是 overwrite 模式且轉換後的檔名與原檔名不同，需將原舊檔移除，避免殘留
        if overwrite and output_path != file_path:
            try:
                os.remove(file_path)
            except Exception as e:
                print(f"[!] 移除舊檔失敗: {e}")
                
        return True
    except Exception as e:
        print(f"[-] 轉換檔案時發生錯誤 {file_path}: {e}")
        # 若發生錯誤，且剛才有做重新命名，則嘗試將備份檔案復原為原名稱
        if not overwrite and os.path.exists(input_path):
            try:
                os.replace(input_path, file_path)
            except:
                pass
        return False

def main():
    parser = argparse.ArgumentParser(description="繁化姬文件批次轉換工具")
    parser.add_argument("path", help="要轉換的檔案路徑或資料夾路徑")
    parser.add_argument(
        "--mode", 
        choices=["Taiwan", "Traditional", "Simplified", "Hongkong", "China"], 
        default="Taiwan", 
        help="轉換模式: Taiwan (台灣化，預設)、Traditional (純簡轉繁)、Simplified (純繁轉簡)、Hongkong (香港化)、China (中國化/大陸用語化)"
    )
    parser.add_argument(
        "--overwrite", 
        action="store_true", 
        help="直接覆蓋原始檔案，而不產生 _old 備份檔（預設會將原簡體檔備份為 *_old）"
    )
    
    args = parser.parse_args()
    
    target_path = os.path.abspath(args.path)
    converter = args.mode
    overwrite = args.overwrite
    
    print(f"[*] 開始執行轉換工作，模式設定: {converter}")
    
    if os.path.isfile(target_path):
        success = process_file(target_path, converter, overwrite)
        if success:
            print("[+] 檔案轉換完成！")
        else:
            print("[-] 檔案轉換失敗。")
    elif os.path.isdir(target_path):
        print(f"[*] 偵測到為目錄，開始遞迴掃描資料夾...")
        supported_exts = [".md", ".txt", ".docx", ".xlsx", ".pptx"]
        files_to_process = []
        
        for root, _, files in os.walk(target_path):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                # 排除已經被我們轉換備份過的舊檔案，避免重複處理
                if ext in supported_exts and "_old" not in file:
                    files_to_process.append(os.path.join(root, file))
                    
        total_files = len(files_to_process)
        print(f"[*] 共發現 {total_files} 個支援的簡體檔案。")
        
        success_count = 0
        for i, file_path in enumerate(files_to_process, 1):
            print(f"[{i}/{total_files}] 正在處理: {os.path.basename(file_path)}")
            if process_file(file_path, converter, overwrite):
                success_count += 1
                
        print(f"[+] 批次處理結束！成功: {success_count}/{total_files} 個檔案。")
    else:
        print(f"[-] 錯誤：不合法的路徑 '{target_path}'")
        sys.exit(1)

if __name__ == "__main__":
    main()
